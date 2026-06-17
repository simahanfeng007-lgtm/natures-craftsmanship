"""L6.72.45 文档解析与主会话安全摘要投影。

该模块只做只读解析与摘要投影：
- 不把 Office/PDF/图片/压缩包等原始字节返回给主会话。
- 可选依赖不可用时返回明确诊断，不把失败误报为 Provider/API 错误。
- human_readable_summary 是唯一适合进入会话气泡的字段。
- L6.72.45 起为追问/引用/导出保留安全 content_preview，不保留原始字节。
"""

from __future__ import annotations

import csv
import io
import json
import re
import zipfile
from html import unescape
from pathlib import Path
from typing import Any, Iterable
import xml.etree.ElementTree as ET

TEXT_ENCODINGS = ("utf-8-sig", "utf-8", "gb18030", "cp936", "utf-16", "utf-16-le", "utf-16-be", "big5")
TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".csv", ".json", ".jsonl", ".html", ".htm", ".xml", ".yaml", ".yml",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".css", ".java", ".cpp", ".c", ".h", ".hpp", ".go", ".rs",
    ".sh", ".bat", ".ps1", ".sql", ".toml", ".ini", ".cfg", ".conf", ".log",
}
DOCUMENT_EXTENSIONS = {".docx", ".pdf", ".xlsx", ".pptx"}
SUPPORTED_PARSE_EXTENSIONS = TEXT_EXTENSIONS | DOCUMENT_EXTENSIONS
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".ico", ".tif", ".tiff", ".svg"}
ARCHIVE_EXTENSIONS = {".zip", ".rar", ".7z", ".tar", ".gz", ".bz2", ".xz"}
BINARY_EXTENSIONS = {".doc", ".xls", ".ppt", ".exe", ".dll", ".so", ".bin", ".dat", ".sqlite", ".db"}


def should_route_to_document_parse(path: str | Path) -> bool:
    """Return whether read_file should delegate to structured document parsing."""
    suffix = Path(str(path or "")).suffix.lower()
    return suffix in SUPPORTED_PARSE_EXTENSIONS or suffix in IMAGE_EXTENSIONS or suffix in ARCHIVE_EXTENSIONS or suffix in BINARY_EXTENSIONS


def _safe_text(value: Any, limit: int = 1000) -> str:
    text = str(value or "")
    text = text.replace("\x00", "")
    text = "".join(ch for ch in text if ch.isprintable() or ch in "\n\r\t")
    text = re.sub(r"\r\n?", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = text.strip()
    if len(text) > limit:
        return text[: max(0, limit - 1)].rstrip() + "…"
    return text


def _clean_lines(lines: Iterable[Any], *, limit: int = 16, line_limit: int = 220) -> list[str]:
    out: list[str] = []
    seen: set[str] = set()
    for item in lines:
        line = _safe_text(item, line_limit)
        if not line:
            continue
        key = line.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append(line)
        if len(out) >= limit:
            break
    return out


def _decode_text_bytes(raw: bytes) -> tuple[str, str, bool]:
    best_text = ""
    best_encoding = ""
    best_score = 10**9
    for encoding in TEXT_ENCODINGS:
        try:
            candidate = raw.decode(encoding)
        except UnicodeDecodeError:
            continue
        replacement = candidate.count("\ufffd")
        controls = sum(1 for ch in candidate if ord(ch) < 32 and ch not in "\n\r\t")
        score = replacement * 25 + controls * 10
        if score < best_score:
            best_text = candidate
            best_encoding = encoding
            best_score = score
            if score == 0:
                break
    if not best_text:
        best_text = raw.decode("utf-8", errors="replace")
        best_encoding = "utf-8-replace"
        best_score = best_text.count("\ufffd") * 25
    text = best_text.replace("\x00", "").replace("\r\n", "\n").replace("\r", "\n")
    uncertain = bool(best_score > max(3, len(text) * 0.02) or text.count("\ufffd") >= 3)
    return text, best_encoding, uncertain


def _read_text_file(path: Path, max_chars: int) -> dict[str, Any]:
    raw = path.read_bytes()
    text, encoding, uncertain = _decode_text_bytes(raw)
    suffix = path.suffix.lower()
    lines = [line.rstrip() for line in text.splitlines()]
    non_empty = [line for line in lines if line.strip()]
    parser = f"text_decoder:{encoding}"
    summary = f"文本文件，约 {len(text)} 字符，{len(lines)} 行。"
    details: list[str] = []
    tables: list[dict[str, Any]] = []

    if suffix == ".json":
        try:
            parsed = json.loads(text)
            if isinstance(parsed, dict):
                details.append("JSON 顶层键：" + ", ".join(_safe_text(k, 60) for k in list(parsed.keys())[:20]))
                summary = f"JSON 对象，顶层键 {len(parsed)} 个。"
            elif isinstance(parsed, list):
                summary = f"JSON 数组，元素 {len(parsed)} 个。"
                details.append("前几个元素类型：" + ", ".join(type(x).__name__ for x in parsed[:8]))
        except Exception as exc:
            details.append(f"JSON 解析诊断：{_safe_text(exc, 160)}")
    elif suffix == ".csv":
        try:
            sample = text[:8192]
            dialect = csv.Sniffer().sniff(sample) if sample.strip() else csv.excel
            reader = csv.reader(io.StringIO(text), dialect)
            rows = []
            for index, row in enumerate(reader):
                rows.append([_safe_text(cell, 80) for cell in row[:12]])
                if index >= 9:
                    break
            if rows:
                summary = f"CSV 表格，预览 {len(rows)} 行，首行 {len(rows[0])} 列。"
                tables.append({"name": path.name, "rows_preview": rows, "row_count_preview": len(rows)})
        except Exception as exc:
            details.append(f"CSV 解析诊断：{_safe_text(exc, 160)}")
    elif suffix in {".html", ".htm"}:
        plain = re.sub(r"<script[\s\S]*?</script>", " ", text, flags=re.I)
        plain = re.sub(r"<style[\s\S]*?</style>", " ", plain, flags=re.I)
        title_match = re.search(r"<title[^>]*>([\s\S]*?)</title>", text, flags=re.I)
        if title_match:
            details.append("HTML 标题：" + _safe_text(unescape(re.sub(r"<[^>]+>", "", title_match.group(1))), 180))
        plain = unescape(re.sub(r"<[^>]+>", " ", plain))
        non_empty = [line for line in re.split(r"\s*\n\s*|(?<=[。.!?？])\s+", plain) if line.strip()] or non_empty
        summary = f"HTML 文档，约 {len(text)} 字符，已抽取可读文本。"

    key_sections = details + _clean_lines(non_empty, limit=14, line_limit=260)
    if uncertain:
        key_sections = ["文本编码不完全确定，已避免输出原始乱码；建议转为 UTF-8 后复核。"] + key_sections
    return {
        "status": "ok" if not uncertain else "partial",
        "parser": parser,
        "summary": summary,
        "key_sections": key_sections[:16],
        "tables": tables,
        "diagnostics": [f"encoding={encoding}", f"bytes={len(raw)}"] + (["encoding_uncertain=true"] if uncertain else []),
        "text_char_count": len(text),
        "content_preview": _safe_text("\n".join(non_empty[:20]), max_chars),
    }


def _xml_text_nodes(xml_bytes: bytes) -> list[str]:
    try:
        root = ET.fromstring(xml_bytes)
    except Exception:
        return []
    parts: list[str] = []
    for node in root.iter():
        if node.tag.endswith("}t") or node.tag.endswith("}title") or node.tag.endswith("}name"):
            if node.text:
                parts.append(node.text)
    return parts


def _parse_docx(path: Path) -> dict[str, Any]:
    paragraphs: list[str] = []
    table_rows: list[list[str]] = []
    parser = "docx_zip_xml"
    diagnostics: list[str] = []
    try:
        from docx import Document  # type: ignore
        parser = "python-docx"
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        for table in doc.tables[:5]:
            for row in table.rows[:8]:
                table_rows.append([_safe_text(cell.text, 100) for cell in row.cells[:8]])
    except Exception as exc:
        diagnostics.append(f"python-docx 不可用或解析失败，已降级 zip+xml：{_safe_text(exc, 160)}")
        try:
            with zipfile.ZipFile(path) as zf:
                if "word/document.xml" in zf.namelist():
                    text_nodes = _xml_text_nodes(zf.read("word/document.xml"))
                    joined: list[str] = []
                    buffer: list[str] = []
                    for node in text_nodes:
                        clean = _safe_text(node, 200)
                        if not clean:
                            continue
                        buffer.append(clean)
                        if clean.endswith(("。", ".", "！", "!", "？", "?")) or len("".join(buffer)) > 140:
                            joined.append("".join(buffer))
                            buffer = []
                    if buffer:
                        joined.append("".join(buffer))
                    paragraphs = joined or text_nodes
        except Exception as zip_exc:
            return {"status": "partial", "parser": parser, "summary": "DOCX 解析失败，可能不是有效 DOCX 或疑似二进制/富文本；未读取原始字节。", "key_sections": ["文件已被安全拦截为文档解析失败诊断，不进入主会话原始字节展示。"], "tables": [], "diagnostics": diagnostics + [_safe_text(zip_exc, 200)]}
    tables = [{"name": "docx_tables", "rows_preview": table_rows[:12], "row_count_preview": len(table_rows)}] if table_rows else []
    summary = f"DOCX 文档，抽取正文段落 {len(paragraphs)} 段" + (f"，表格预览 {len(table_rows)} 行。" if table_rows else "。")
    content_preview = _safe_text("\n".join(paragraphs[:80]), 12000)
    return {"status": "ok" if paragraphs or table_rows else "partial", "parser": parser, "summary": summary, "key_sections": _clean_lines(paragraphs, limit=16, line_limit=260), "tables": tables, "diagnostics": diagnostics, "content_preview": content_preview}


def _parse_xlsx(path: Path) -> dict[str, Any]:
    diagnostics: list[str] = []
    try:
        import openpyxl  # type: ignore
    except Exception as exc:
        return {
            "status": "partial",
            "parser": "xlsx_dependency_missing",
            "summary": "XLSX 文件已识别，但当前环境缺少 openpyxl，未读取原始字节。",
            "key_sections": ["需要安装 openpyxl 后才能提取工作表、表头和单元格摘要。"],
            "sheets": [],
            "diagnostics": [f"openpyxl_import_error={_safe_text(exc, 160)}"],
        }
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        sheets: list[dict[str, Any]] = []
        key_sections: list[str] = []
        for ws in wb.worksheets[:8]:
            rows_preview: list[list[str]] = []
            for r_idx, row in enumerate(ws.iter_rows(max_row=12, max_col=10, values_only=True), start=1):
                values = [_safe_text(cell, 80) for cell in row]
                if any(values):
                    rows_preview.append(values)
                if r_idx >= 12:
                    break
            headers = rows_preview[0] if rows_preview else []
            sheets.append({
                "name": ws.title,
                "max_row": ws.max_row,
                "max_column": ws.max_column,
                "headers": headers,
                "rows_preview": rows_preview[:8],
            })
            key_sections.append(f"工作表 {ws.title}：约 {ws.max_row} 行 × {ws.max_column} 列；表头：{', '.join([x for x in headers if x][:8]) or '未识别'}")
        content_lines = list(key_sections)
        for sheet in sheets[:8]:
            content_lines.append(f"工作表 {sheet.get('name')} 预览：")
            for row in sheet.get("rows_preview", [])[:8]:
                content_lines.append(" | ".join(_safe_text(cell, 90) for cell in row))
        return {
            "status": "ok",
            "parser": "openpyxl",
            "summary": f"XLSX 工作簿，工作表 {len(wb.worksheets)} 个。",
            "key_sections": key_sections,
            "sheets": sheets,
            "tables": [{"name": s["name"], "rows_preview": s.get("rows_preview", [])} for s in sheets[:4]],
            "diagnostics": diagnostics,
            "content_preview": _safe_text("\n".join(content_lines), 12000),
        }
    except Exception as exc:
        return {"status": "partial", "parser": "openpyxl", "summary": "XLSX 解析失败，可能不是有效 XLSX 或疑似二进制/富文本；未读取原始字节。", "key_sections": ["文件已被安全拦截为文档解析失败诊断，不进入主会话原始字节展示。"], "sheets": [], "diagnostics": [_safe_text(exc, 200)]}


def _parse_pptx(path: Path) -> dict[str, Any]:
    slides: list[dict[str, Any]] = []
    diagnostics: list[str] = []
    parser = "pptx_zip_xml"
    try:
        from pptx import Presentation  # type: ignore
        parser = "python-pptx"
        prs = Presentation(str(path))
        for idx, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                value = getattr(shape, "text", "")
                if value and str(value).strip():
                    texts.append(_safe_text(value, 300))
            slides.append({"index": idx, "text": _clean_lines(texts, limit=10, line_limit=220)})
    except Exception as exc:
        diagnostics.append(f"python-pptx 不可用或解析失败，已降级 zip+xml：{_safe_text(exc, 160)}")
        try:
            with zipfile.ZipFile(path) as zf:
                names = sorted(name for name in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", name))
                for idx, name in enumerate(names[:80], start=1):
                    texts = _clean_lines(_xml_text_nodes(zf.read(name)), limit=12, line_limit=220)
                    slides.append({"index": idx, "text": texts})
        except Exception as zip_exc:
            return {"status": "partial", "parser": parser, "summary": "PPTX 解析失败，可能不是有效 PPTX 或疑似二进制/富文本；未读取原始字节。", "key_sections": ["文件已被安全拦截为文档解析失败诊断，不进入主会话原始字节展示。"], "slides": [], "diagnostics": diagnostics + [_safe_text(zip_exc, 200)]}
    key_sections = [f"第 {slide['index']} 页：" + ("；".join(slide.get("text", [])[:4]) or "未抽取到文本") for slide in slides[:20]]
    content_preview = _safe_text("\n".join(f"第 {slide.get('index')} 页：" + "；".join(slide.get("text", [])) for slide in slides[:80]), 12000)
    return {"status": "ok" if slides else "partial", "parser": parser, "summary": f"PPTX 演示文稿，抽取幻灯片 {len(slides)} 页。", "key_sections": key_sections, "slides": slides, "diagnostics": diagnostics, "content_preview": content_preview}


def _parse_pdf(path: Path) -> dict[str, Any]:
    diagnostics: list[str] = []
    pages: list[str] = []
    try:
        from pypdf import PdfReader  # type: ignore
        reader = PdfReader(str(path))
        for page_index, page in enumerate(reader.pages[:12], start=1):
            try:
                pages.append(_safe_text(page.extract_text() or "", 1200))
            except Exception as exc:
                diagnostics.append(f"page_extract_error={page_index}:{_safe_text(exc, 120)}")
        page_items = [{"index": idx, "text": text} for idx, text in enumerate(pages, start=1) if text]
        return {
            "status": "ok" if any(pages) else "partial",
            "parser": "pypdf",
            "summary": f"PDF 文档，页数 {len(reader.pages)}，已抽取前 {min(len(reader.pages), 12)} 页文本。",
            "key_sections": _clean_lines(pages, limit=12, line_limit=300),
            "pages": page_items,
            "content_preview": _safe_text("\n".join(pages), 12000),
            "diagnostics": diagnostics,
        }
    except Exception as exc:
        diagnostics.append(f"pypdf 不可用或解析失败：{_safe_text(exc, 160)}")
    try:
        import pdfplumber  # type: ignore
        with pdfplumber.open(str(path)) as pdf:
            for page_index, page in enumerate(pdf.pages[:12], start=1):
                pages.append(_safe_text(page.extract_text() or "", 1200))
            page_items = [{"index": idx, "text": text} for idx, text in enumerate(pages, start=1) if text]
            return {
                "status": "ok" if any(pages) else "partial",
                "parser": "pdfplumber",
                "summary": f"PDF 文档，页数 {len(pdf.pages)}，已抽取前 {min(len(pdf.pages), 12)} 页文本。",
                "key_sections": _clean_lines(pages, limit=12, line_limit=300),
                "pages": page_items,
                "content_preview": _safe_text("\n".join(pages), 12000),
                "diagnostics": diagnostics,
            }
    except Exception as exc:
        diagnostics.append(f"pdfplumber 不可用或解析失败：{_safe_text(exc, 160)}")
    return {
        "status": "partial",
        "parser": "pdf_dependency_missing",
        "summary": "PDF 文件已识别，但当前环境缺少可用 PDF 文本解析依赖或文件不可抽取；未读取原始字节。",
        "key_sections": ["需要安装 pypdf 或 pdfplumber 后再抽取 PDF 正文；当前只返回文件元信息与明确诊断。"],
        "diagnostics": diagnostics,
    }


def _classify_unparsed(path: Path) -> dict[str, Any]:
    suffix = path.suffix.lower()
    if suffix in IMAGE_EXTENSIONS:
        file_type = "image"
        suggestion = "图片文件不进入 read_file 原始字节展示；如需视觉理解，请使用图片/视觉解析工具。"
    elif suffix in ARCHIVE_EXTENSIONS:
        file_type = "archive"
        suggestion = "压缩包不进入主会话原始字节展示；如需处理，请先列清单或使用解包/压缩包解析工具。"
    elif suffix in BINARY_EXTENSIONS:
        file_type = "binary_or_legacy_office"
        suggestion = "二进制/旧 Office 文件不直接展示原始字节；建议转换为 docx/xlsx/pptx/pdf 或使用专用解析器。"
    else:
        file_type = "binary_or_unknown"
        suggestion = "未知二进制文件不直接展示原始字节；请指定合适的解析工具或转换为文本格式。"
    return {
        "status": "partial",
        "parser": "binary_guard",
        "file_type": file_type,
        "summary": f"{file_type} 文件已识别；主会话仅显示元信息，不读取原始字节。",
        "key_sections": [suggestion],
        "diagnostics": ["raw_bytes_hidden=true"],
    }


def _file_type_for_suffix(suffix: str) -> str:
    if suffix in TEXT_EXTENSIONS:
        return "text/code"
    if suffix == ".docx":
        return "docx"
    if suffix == ".xlsx":
        return "xlsx"
    if suffix == ".pptx":
        return "pptx"
    if suffix == ".pdf":
        return "pdf"
    if suffix in IMAGE_EXTENSIONS:
        return "image"
    if suffix in ARCHIVE_EXTENSIONS:
        return "archive"
    if suffix in BINARY_EXTENSIONS:
        return "binary"
    return "unknown"


def _format_table_preview(tables: list[dict[str, Any]] | None, sheets: list[dict[str, Any]] | None, slides: list[dict[str, Any]] | None) -> list[str]:
    out: list[str] = []
    for table in (tables or [])[:3]:
        rows = table.get("rows_preview") if isinstance(table, dict) else []
        if rows:
            first = rows[0]
            out.append(f"表格 {table.get('name', 'table')}：预览 {len(rows)} 行；首行：{', '.join(_safe_text(x, 60) for x in first[:8])}")
    for sheet in (sheets or [])[:4]:
        out.append(f"Sheet {sheet.get('name')}：{sheet.get('max_row', '?')} 行 × {sheet.get('max_column', '?')} 列。")
    for slide in (slides or [])[:6]:
        text = "；".join(slide.get("text", [])[:3]) if isinstance(slide, dict) else ""
        out.append(f"幻灯片 {slide.get('index', '?')}：{text or '未抽取到文本'}")
    return out


def build_human_readable_summary(result: dict[str, Any]) -> str:
    lines: list[str] = ["【文档解析】"]
    lines.append(f"- 文件名：{_safe_text(result.get('file_name'), 180)}")
    lines.append(f"- 文件类型：{_safe_text(result.get('file_type'), 80)}")
    lines.append(f"- 解析方式：{_safe_text(result.get('parser'), 120)}")
    lines.append(f"- 大小：{result.get('size_bytes', 0)} bytes")
    lines.append(f"- 状态：{_safe_text(result.get('status'), 40)}")
    lines.append(f"- 摘要：{_safe_text(result.get('summary'), 300)}")
    key_sections = _clean_lines(result.get("key_sections") or [], limit=10, line_limit=260)
    if key_sections:
        lines.append("- 关键段落 / 内容摘要：")
        for idx, item in enumerate(key_sections[:10], start=1):
            lines.append(f"  {idx}. {item}")
    table_lines = _format_table_preview(result.get("tables"), result.get("sheets"), result.get("slides"))
    if table_lines:
        lines.append("- 表格 / 工作表 / 幻灯片摘要：")
        for item in table_lines[:10]:
            lines.append(f"  - {item}")
    suggestions = result.get("next_suggestions") or []
    if suggestions:
        lines.append("- 后续建议：")
        for item in suggestions[:4]:
            lines.append(f"  - {_safe_text(item, 180)}")
    return "\n".join(lines)


def parse_document(path: str | Path, *, max_chars: int = 8000) -> dict[str, Any]:
    target = Path(path).expanduser()
    result: dict[str, Any]
    if not target.exists():
        result = {"status": "failed", "parser": "path_guard", "summary": "文件不存在。", "key_sections": [], "diagnostics": ["path_not_found"]}
    elif not target.is_file():
        result = {"status": "failed", "parser": "path_guard", "summary": "目标不是文件。", "key_sections": [], "diagnostics": ["not_file"]}
    else:
        suffix = target.suffix.lower()
        try:
            if suffix in TEXT_EXTENSIONS:
                result = _read_text_file(target, max_chars)
            elif suffix == ".docx":
                result = _parse_docx(target)
            elif suffix == ".xlsx":
                result = _parse_xlsx(target)
            elif suffix == ".pptx":
                result = _parse_pptx(target)
            elif suffix == ".pdf":
                result = _parse_pdf(target)
            else:
                result = _classify_unparsed(target)
        except OSError as exc:
            result = {"status": "failed", "parser": "os_error", "summary": f"文件读取失败：{_safe_text(exc, 200)}", "key_sections": [], "diagnostics": ["os_error"]}
        except Exception as exc:
            result = {"status": "partial", "parser": "document_parse_exception", "summary": f"文档解析失败：{_safe_text(exc, 200)}；疑似二进制/富文本，未读取原始字节。", "key_sections": ["文件解析失败已转为明确诊断，不误报 Provider/API 错误。"], "diagnostics": [type(exc).__name__]}
    try:
        size_bytes = target.stat().st_size if target.exists() and target.is_file() else 0
    except Exception:
        size_bytes = 0
    result.setdefault("status", "partial")
    result.setdefault("parser", "document_parse")
    result.setdefault("summary", "已完成文档解析投影。")
    result.setdefault("key_sections", [])
    result.setdefault("diagnostics", [])
    result.setdefault("tables", [])
    result.setdefault("sheets", [])
    result.setdefault("slides", [])
    result.setdefault("pages", [])
    result.update({
        "file_name": target.name,
        "file_path": str(target),
        "file_type": _file_type_for_suffix(target.suffix.lower()),
        "suffix": target.suffix.lower(),
        "size_bytes": size_bytes,
        "raw_bytes_hidden": True,
        "main_chat_projection_only": True,
    })
    result["next_suggestions"] = [
        "可以继续追问：指定关键词、页码、工作表、幻灯片编号或引用编号。",
        "可以导出：使用 document_export 生成 md/txt/json 摘要与引用片段。",
        "需要修改原文时进入长链模式，先生成 document_rewrite_plan；明确 old_text/new_text 后用 document_apply_rewrite 生成修订副本/写回。",
        "解析失败时补装对应依赖，或转换为 txt/md/csv 后重试。",
    ]
    result["human_readable_summary"] = build_human_readable_summary(result)
    return result
