"""L6.72.46 文档真实修改写回与回滚 adapters。

边界：
- 只在 RuntimeToolRegistry / ExecutionSpine / QualityGate / Audit 链内执行。
- 默认生成修订副本；只有 overwrite=true/mode=overwrite 才覆盖原文件，并在覆盖前写入备份。
- 支持 txt/md/csv/json/html/code/docx/xlsx/pptx 的受控 literal replace / full text 写入。
- PDF 不做原文写回；返回明确诊断，建议导出修订草案。
- 主会话只输出摘要、路径、命中数、operation_id，不输出原始正文/raw bytes/stderr。
"""

from __future__ import annotations

import hashlib
import json
import re
import shutil
from datetime import datetime
from pathlib import Path
from typing import Any

from ..document_context_store import load_document_context, safe_text
from ..document_parser import TEXT_EXTENSIONS, parse_document
from ..host_path_normalizer import normalize_argument_path, normalization_public_data
from ..physical_commit import verify_file_commit, write_text_atomic_verified
from ..result_normalizer import truncate_text
from ..tool_invocation import ToolInvocation
from ..tool_result import ToolResult, ToolResultStatus
from ..turn_context import TurnContext
from ..workspace_guard import WorkspaceGuard, WorkspaceViolation

WRITEBACK_SCHEMA = "tiangong.l6_72_46.document_writeback_manifest.v1"
ROLLBACK_SCHEMA = "tiangong.l6_72_46.document_rollback_manifest.v1"
_WRITEBACK_DIR = Path(".linyuanzhe") / "document_writeback"
_SUPPORTED_OFFICE_SUFFIXES = {".docx", ".xlsx", ".pptx"}
_SUPPORTED_TEXT_SUFFIXES = set(TEXT_EXTENSIONS)
_DIRECT_WRITE_SUFFIXES = _SUPPORTED_TEXT_SUFFIXES | _SUPPORTED_OFFICE_SUFFIXES
_PDF_SUFFIXES = {".pdf"}
_TEXT_ENCODINGS = ("utf-8-sig", "utf-8", "gb18030", "cp936", "utf-16", "utf-16-le", "utf-16-be", "big5")

_WINDOWS_PROTECTED_ROOTS = {
    "windows",
    "program files",
    "program files (x86)",
    "programdata",
    "system volume information",
    "recovery",
    "$recycle.bin",
}


def _bool(value: Any, default: bool = False) -> bool:
    if value in (None, ""):
        return default
    if isinstance(value, bool):
        return value
    text = str(value).strip().lower()
    if text in {"1", "true", "yes", "y", "on", "是", "开启", "覆盖", "overwrite"}:
        return True
    if text in {"0", "false", "no", "n", "off", "否", "关闭"}:
        return False
    return default


def _sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def _operation_id(source: Path, target: Path, instruction: str) -> str:
    payload = f"{source}|{target}|{instruction}|{datetime.now().isoformat(timespec='microseconds')}"
    return "docwrite_" + hashlib.sha256(payload.encode("utf-8", errors="ignore")).hexdigest()[:16]


def _safe_rel(path: Path, workspace: Path) -> str:
    try:
        return path.resolve().relative_to(workspace.resolve()).as_posix()
    except Exception:
        return str(path)


def _manifest_root(workspace: Path) -> Path:
    root = workspace / _WRITEBACK_DIR
    root.mkdir(parents=True, exist_ok=True)
    (root / "backups").mkdir(parents=True, exist_ok=True)
    (root / "manifests").mkdir(parents=True, exist_ok=True)
    return root


def _manifest_path(workspace: Path, operation_id: str) -> Path:
    return _manifest_root(workspace) / "manifests" / f"{safe_text(operation_id, 80)}.json"


def _backup_file(workspace: Path, path: Path, operation_id: str, label: str) -> str:
    if not path.exists():
        return ""
    backup_dir = _manifest_root(workspace) / "backups" / safe_text(operation_id, 80)
    backup_dir.mkdir(parents=True, exist_ok=True)
    backup = backup_dir / f"{label}_{path.name}"
    # 避免同名覆盖。
    if backup.exists():
        stem = backup.stem
        suffix = backup.suffix
        for idx in range(2, 100):
            candidate = backup.with_name(f"{stem}_{idx}{suffix}")
            if not candidate.exists():
                backup = candidate
                break
    shutil.copy2(path, backup)
    return _safe_rel(backup, workspace)


def _looks_like_windows_protected_target(target: Path, workspace: Path) -> bool:
    try:
        rel_parts = [part.lower() for part in target.resolve().relative_to(workspace.resolve()).parts]
    except Exception:
        rel_parts = [part.lower() for part in target.parts]
    if not rel_parts:
        return False
    if rel_parts[0] in _WINDOWS_PROTECTED_ROOTS:
        return True
    return "system32" in rel_parts or "syswow64" in rel_parts


def _decode_text_file(path: Path) -> tuple[str, str, bool]:
    raw = path.read_bytes()
    best = ""
    best_encoding = "utf-8"
    best_score = 10**9
    for encoding in _TEXT_ENCODINGS:
        try:
            candidate = raw.decode(encoding)
        except UnicodeDecodeError:
            continue
        score = candidate.count("\ufffd") * 25 + sum(1 for ch in candidate if ord(ch) < 32 and ch not in "\n\r\t") * 10
        if score < best_score:
            best = candidate
            best_encoding = encoding
            best_score = score
            if score == 0:
                break
    if not best:
        best = raw.decode("utf-8", errors="replace")
        best_encoding = "utf-8"
        best_score = best.count("\ufffd") * 25
    text = best.replace("\x00", "").replace("\r\n", "\n").replace("\r", "\n")
    uncertain = bool(best_score > max(3, len(text) * 0.02) or text.count("\ufffd") >= 3)
    return text, best_encoding, uncertain


def _replacement_items(args: dict[str, Any], instruction: str) -> list[dict[str, str]]:
    items: list[dict[str, str]] = []
    raw = args.get("replacements") or args.get("replace") or args.get("edits")
    if isinstance(raw, list):
        for item in raw:
            if not isinstance(item, dict):
                continue
            old = item.get("old_text") or item.get("old") or item.get("from") or item.get("find") or item.get("原文")
            new = item.get("new_text") or item.get("new") or item.get("to") or item.get("replace_with") or item.get("替换为")
            if old not in (None, "") and new not in (None, ""):
                items.append({"old_text": safe_text(old, 500), "new_text": safe_text(new, 1000)})
    elif isinstance(raw, dict):
        for old, new in raw.items():
            if old not in (None, "") and new not in (None, ""):
                items.append({"old_text": safe_text(old, 500), "new_text": safe_text(new, 1000)})

    old_text = args.get("old_text") or args.get("old") or args.get("find") or args.get("from") or args.get("原文")
    new_text = args.get("new_text") or args.get("new") or args.get("replace_with") or args.get("to") or args.get("替换为")
    if old_text not in (None, "") and new_text not in (None, ""):
        items.append({"old_text": safe_text(old_text, 500), "new_text": safe_text(new_text, 1000)})

    if not items:
        items.extend(_parse_instruction_replacements(instruction))

    clean: list[dict[str, str]] = []
    seen: set[tuple[str, str]] = set()
    for item in items:
        old = _strip_quote_noise(item.get("old_text") or "")
        new = _strip_quote_noise(item.get("new_text") or "")
        if len(old) < 2:
            continue
        key = (old, new)
        if key in seen:
            continue
        seen.add(key)
        clean.append({"old_text": old, "new_text": new})
        if len(clean) >= 20:
            break
    return clean


def _strip_quote_noise(text: str) -> str:
    value = safe_text(text, 1000).strip()
    value = value.strip("\"'“”‘’《》<>` ")
    value = re.sub(r"^.*\.(?:docx|pdf|xlsx|xlsm|pptx|csv|md|markdown|txt|html?|json|py|js|ts|tsx|jsx|css|java|cpp|c|go|rs)(?:\s*(?:里|里的|中|中的|内|里面|的))?", "", value, flags=re.IGNORECASE).strip()
    value = re.sub(r"^(?:帮我|请|把)?\s*(?:桌面|下载|文档|我的文档|desktop|downloads?|documents?)(?:的|里|里的|中|中的)?\s*", "", value, flags=re.IGNORECASE).strip()
    value = re.sub(r"^(?:这个|该|刚才|上面|这份)?(?:文档|原文|正文|内容|文件|资料)(?:中|里|里的|的)?", "", value).strip()
    value = re.sub(r"^(?:里面|里|中|的)", "", value).strip()
    value = re.sub(r"(?:并)?(?:写回|保存修改|应用修改|生成修订副本|覆盖写入|真正修改|执行修改)$", "", value).strip()
    # L6.72.59：代码/配置类 literal replacement 可能需要保留行尾冒号、分号等语法字符。
    # 旧逻辑会把 new_text="def f():" 清洗成 "def f()"，导致工具报告写回成功但文件实际未修复。
    return value.strip("，,。. ")


def _parse_instruction_replacements(instruction: str) -> list[dict[str, str]]:
    text = safe_text(instruction, 1600)
    patterns = (
        r"把\s*[“\"'《]?(?P<old>[^“”\"'《》\n]{2,180})[”\"'》]?\s*(?:替换成|替换为|改成|改为|换成)\s*[“\"'《]?(?P<new>[^“”\"'《》\n]{0,260})[”\"'》]?",
        r"(?:将|把)\s*(?P<old>[^\n]{2,180})\s*(?:改成|改为|替换成|替换为)\s*(?P<new>[^\n]{0,260})",
        r"replace\s+[\"']?(?P<old>[^\"'\n]{2,180})[\"']?\s+with\s+[\"']?(?P<new>[^\"'\n]{0,260})[\"']?",
    )
    out: list[dict[str, str]] = []
    for pattern in patterns:
        for match in re.finditer(pattern, text, flags=re.IGNORECASE):
            old = _strip_quote_noise(match.group("old"))
            new = _strip_quote_noise(match.group("new"))
            if old:
                out.append({"old_text": old, "new_text": new})
    return out


def _apply_replacements_to_text(text: str, replacements: list[dict[str, str]]) -> tuple[str, int, list[dict[str, Any]]]:
    changed = text
    total = 0
    details: list[dict[str, Any]] = []
    for item in replacements:
        old = item["old_text"]
        new = item["new_text"]
        count = changed.count(old)
        if count:
            changed = changed.replace(old, new)
        total += count
        details.append({"old_text": safe_text(old, 160), "new_text": safe_text(new, 160), "count": count})
    return changed, total, details


def _apply_replacements_to_docx(source: Path, target: Path, replacements: list[dict[str, str]]) -> tuple[int, list[dict[str, Any]]]:
    try:
        from docx import Document  # type: ignore
    except Exception as exc:
        raise RuntimeError(f"python-docx 不可用，无法写回 DOCX：{safe_text(exc, 160)}") from exc
    doc = Document(str(source))
    total = 0
    details: list[dict[str, Any]] = []

    def patch_paragraph(paragraph: Any) -> int:
        local = 0
        for item in replacements:
            old = item["old_text"]
            new = item["new_text"]
            run_hits = 0
            for run in paragraph.runs:
                if old in run.text:
                    hits = run.text.count(old)
                    run.text = run.text.replace(old, new)
                    run_hits += hits
            if run_hits:
                local += run_hits
                continue
            if old in paragraph.text:
                hits = paragraph.text.count(old)
                paragraph.text = paragraph.text.replace(old, new)
                local += hits
        return local

    for p in doc.paragraphs:
        total += patch_paragraph(p)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    total += patch_paragraph(p)
    for item in replacements:
        details.append({"old_text": safe_text(item["old_text"], 160), "new_text": safe_text(item["new_text"], 160), "count": "see_total_docx"})
    if total > 0:
        target.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(target))
    return total, details


def _apply_replacements_to_xlsx(source: Path, target: Path, replacements: list[dict[str, str]]) -> tuple[int, list[dict[str, Any]]]:
    try:
        import openpyxl  # type: ignore
    except Exception as exc:
        raise RuntimeError(f"openpyxl 不可用，无法写回 XLSX：{safe_text(exc, 160)}") from exc
    wb = openpyxl.load_workbook(str(source))
    total = 0
    details: list[dict[str, Any]] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if not isinstance(cell.value, str):
                    continue
                text = cell.value
                changed, count, _ = _apply_replacements_to_text(text, replacements)
                if count:
                    cell.value = changed
                    total += count
    for item in replacements:
        details.append({"old_text": safe_text(item["old_text"], 160), "new_text": safe_text(item["new_text"], 160), "count": "see_total_xlsx"})
    if total > 0:
        target.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(target))
    return total, details


def _apply_replacements_to_pptx(source: Path, target: Path, replacements: list[dict[str, str]]) -> tuple[int, list[dict[str, Any]]]:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:
        raise RuntimeError(f"python-pptx 不可用，无法写回 PPTX：{safe_text(exc, 160)}") from exc
    prs = Presentation(str(source))
    total = 0
    details: list[dict[str, Any]] = []

    def patch_text_frame(tf: Any) -> int:
        local = 0
        for paragraph in tf.paragraphs:
            for item in replacements:
                old = item["old_text"]
                new = item["new_text"]
                run_hits = 0
                for run in paragraph.runs:
                    if old in run.text:
                        hits = run.text.count(old)
                        run.text = run.text.replace(old, new)
                        run_hits += hits
                if run_hits:
                    local += run_hits
        return local

    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                total += patch_text_frame(shape.text_frame)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if getattr(cell, "text_frame", None) is not None:
                            total += patch_text_frame(cell.text_frame)
    for item in replacements:
        details.append({"old_text": safe_text(item["old_text"], 160), "new_text": safe_text(item["new_text"], 160), "count": "see_total_pptx"})
    if total > 0:
        target.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(target))
    return total, details


def _resolve_context(invocation: ToolInvocation, context: TurnContext) -> tuple[dict[str, Any] | None, str]:
    args = invocation.arguments
    document_id = safe_text(args.get("document_id"), 100)
    raw_path = args.get("path") or args.get("source") or ""
    if document_id:
        ctx = load_document_context(context.workspace, document_id=document_id)
        return ctx, "loaded_from_document_id" if ctx else "missing_context"
    if raw_path:
        # 先尝试加载已有 context；不存在时只读解析并建立 context，避免写回工具依赖前置手工解析。
        guard = WorkspaceGuard(context.workspace)
        target = guard.resolve_for_read(raw_path)
        ctx = load_document_context(context.workspace, path=target)
        if ctx:
            return ctx, "loaded_from_path_context"
        parsed = parse_document(target, max_chars=int(args.get("max_chars") or context.policy.max_output_chars))
        from ..document_context_store import save_document_context  # 延迟导入，避免循环意图。
        ctx, _ = save_document_context(context.workspace, parsed)
        return ctx, "parsed_from_path"
    ctx = load_document_context(context.workspace)
    return ctx, "loaded_last_context" if ctx else "missing_context"


def _source_path_from_context_or_args(invocation: ToolInvocation, context: TurnContext, ctx: dict[str, Any] | None) -> str:
    args = invocation.arguments
    direct = args.get("path") or args.get("source") or args.get("file_path")
    if direct:
        return str(direct)
    if ctx:
        meta = dict(ctx.get("metadata") or {})
        if meta.get("file_path"):
            return str(meta.get("file_path"))
    return ""


def _default_target_for(source: Path) -> Path:
    return source.with_name(f"{source.stem}.linyuanzhe_rewrite{source.suffix}")


def _write_manifest(workspace: Path, operation_id: str, payload: dict[str, Any]) -> str:
    path = _manifest_path(workspace, operation_id)
    commit = write_text_atomic_verified(path, json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    if not commit.get("physical_commit_verified"):
        raise OSError("writeback manifest physical commit verification failed")
    return _safe_rel(path, workspace)


def document_apply_rewrite_adapter(invocation: ToolInvocation, context: TurnContext) -> ToolResult:
    """Apply a governed document rewrite to a copy or overwritable target."""
    guard = WorkspaceGuard(context.workspace)
    try:
        ctx, context_source = _resolve_context(invocation, context)
        source_hint = _source_path_from_context_or_args(invocation, context, ctx)
        if not source_hint:
            return ToolResult(
                invocation.step_id,
                invocation.tool_name,
                ToolResultStatus.FAILED,
                "没有可写回的文档源路径。请先 document_parse，或在 document_apply_rewrite 中传入 path。",
                error_code="document_source_missing",
            )
        normalized_source_hint, source_path_normalization = normalize_argument_path(source_hint, context.user_message)
        source = guard.resolve_for_read(normalized_source_hint)
        if not source.exists() or not source.is_file():
            return ToolResult(invocation.step_id, invocation.tool_name, ToolResultStatus.FAILED, "文档源文件不存在或不是文件。", error_code="source_not_file")
        suffix = source.suffix.lower()
        if suffix in _PDF_SUFFIXES:
            return ToolResult(
                invocation.step_id,
                invocation.tool_name,
                ToolResultStatus.BLOCKED,
                "PDF 不支持安全原文写回。请使用 document_export 生成 md/txt 修订草案，或转换为 docx 后再写回；未读取或写入 PDF 原始字节。",
                error_code="pdf_writeback_unsupported",
                data={"source": _safe_rel(source, context.workspace), "raw_bytes_hidden": True},
            )
        if suffix not in _DIRECT_WRITE_SUFFIXES:
            return ToolResult(
                invocation.step_id,
                invocation.tool_name,
                ToolResultStatus.BLOCKED,
                f"当前文件类型 {suffix or '<无后缀>'} 不支持直接写回；已拒绝读取/写入原始二进制。",
                error_code="unsupported_writeback_type",
                data={"source": _safe_rel(source, context.workspace), "raw_bytes_hidden": True},
            )

        args = dict(invocation.arguments)
        instruction = safe_text(args.get("instruction") or args.get("query") or context.user_message, 1600)
        replacements = _replacement_items(args, instruction)
        content = args.get("content") if args.get("content") not in (None, "") else args.get("new_content")
        operation = safe_text(args.get("operation") or ("replace" if replacements else ("write_full_text" if content not in (None, "") and suffix in _SUPPORTED_TEXT_SUFFIXES else "")), 80)
        overwrite = _bool(args.get("overwrite"), default=False) or safe_text(args.get("mode"), 40).lower() == "overwrite"
        dry_run = _bool(args.get("dry_run"), default=False)
        allow_no_match = _bool(args.get("allow_no_match"), default=False)
        target_hint = args.get("target") or args.get("output_path")
        if target_hint:
            normalized_target_hint, target_path_normalization = normalize_argument_path(target_hint, context.user_message)
            target = guard.resolve_for_write(normalized_target_hint)
        else:
            target_path_normalization = source_path_normalization
            target = source if overwrite else _default_target_for(source)
        if target.resolve() == source.resolve() and not overwrite:
            target = _default_target_for(source)
        if _looks_like_windows_protected_target(target, context.workspace):
            return ToolResult(
                invocation.step_id,
                invocation.tool_name,
                ToolResultStatus.BLOCKED,
                "Windows 管理员权限目录写回需要显式提权；已拒绝直接写入。",
                error_code="windows_permission_required",
                data={"target": str(target), "requires_administrator": True},
            )
        if not replacements and content in (None, ""):
            # Q25 自愈式错误：读出文件原文前 800 字，让 LLM 看到后能补全参数重试。
            try:
                preview_text, _enc, _enc_uncertain = _decode_text_file(source)
            except Exception:
                preview_text = ""
            preview = (preview_text or "")[:800]
            hint = f"\n文件当前内容预览（前800字）：\n---\n{preview}\n---\n请根据以上内容指定 old_text/new_text、replacements，或用 content + operation:append 追加。"
            return ToolResult(
                invocation.step_id,
                invocation.tool_name,
                ToolResultStatus.FAILED,
                f"写回缺少明确修改单元。{hint}",
                error_code="document_rewrite_instruction_ambiguous",
                data={"context_source": context_source, "source": _safe_rel(source, context.workspace), "content_preview_chars": len(preview), "raw_bytes_hidden": True},
            )

        op_id = _operation_id(source, target, instruction)
        source_hash = _sha256_file(source)
        applied = 0
        details: list[dict[str, Any]] = []
        encoding = ""
        encoding_uncertain = False
        if suffix in _SUPPORTED_TEXT_SUFFIXES:
            current_text, encoding, encoding_uncertain = _decode_text_file(source)
            if content not in (None, "") and operation in {"write_full_text", "full_text", "replace_all"}:
                new_text = str(content)
                applied = 1 if new_text != current_text else 0
                details = [{"operation": "write_full_text", "count": applied, "chars": len(new_text)}]
            elif content not in (None, "") and operation in {"append", "append_text"}:
                new_text = current_text.rstrip() + "\n" + str(content).strip() + "\n"
                applied = 1
                details = [{"operation": "append", "count": 1, "chars": len(str(content))}]
            else:
                new_text, applied, details = _apply_replacements_to_text(current_text, replacements)
            if applied <= 0 and not allow_no_match:
                return ToolResult(
                    invocation.step_id,
                    invocation.tool_name,
                    ToolResultStatus.FAILED,
                    "未命中可替换文本，已停止写回。请先用 document_query 定位引用片段，或设置 allow_no_match=true。",
                    error_code="document_rewrite_no_match",
                    data={"source": _safe_rel(source, context.workspace), "replacements": details, "raw_bytes_hidden": True},
                )
            if not dry_run:
                target.parent.mkdir(parents=True, exist_ok=True)
                backup_rel = _backup_file(context.workspace, target, op_id, "target_before") if target.exists() else ""
                commit = write_text_atomic_verified(target, new_text, encoding="utf-8" if encoding_uncertain else (encoding or "utf-8"))
                if not commit.get("physical_commit_verified"):
                    return ToolResult(
                        invocation.step_id,
                        invocation.tool_name,
                        ToolResultStatus.FAILED,
                        "文档写回后物理落盘验真失败：目标文件未通过 read-after-write 校验。",
                        error_code="physical_commit_verification_failed",
                        data={"source": _safe_rel(source, context.workspace), "target": _safe_rel(target, context.workspace), "commit": commit, "raw_bytes_hidden": True},
                    )
            else:
                backup_rel = ""
        else:
            if not replacements:
                return ToolResult(
                    invocation.step_id,
                    invocation.tool_name,
                    ToolResultStatus.FAILED,
                    "Office 文档写回需要明确 old_text/new_text 或 replacements；未写入任何文件。",
                    error_code="office_rewrite_requires_replacements",
                    data={"source": _safe_rel(source, context.workspace), "raw_bytes_hidden": True},
                )
            if dry_run:
                backup_rel = ""
                parsed = parse_document(source, max_chars=context.policy.max_output_chars)
                preview_text = str(parsed.get("content_preview") or "")
                _, applied, details = _apply_replacements_to_text(preview_text, replacements)
            else:
                target.parent.mkdir(parents=True, exist_ok=True)
                backup_rel = _backup_file(context.workspace, target, op_id, "target_before") if target.exists() else ""
                if suffix == ".docx":
                    applied, details = _apply_replacements_to_docx(source, target, replacements)
                elif suffix == ".xlsx":
                    applied, details = _apply_replacements_to_xlsx(source, target, replacements)
                elif suffix == ".pptx":
                    applied, details = _apply_replacements_to_pptx(source, target, replacements)
                else:
                    applied = 0
                if applied > 0:
                    office_commit = verify_file_commit(target)
                    if not office_commit.get("physical_commit_verified"):
                        return ToolResult(
                            invocation.step_id,
                            invocation.tool_name,
                            ToolResultStatus.FAILED,
                            "Office 文档写回后物理落盘验真失败：目标文件未通过存在性/父目录枚举校验。",
                            error_code="physical_commit_verification_failed",
                            data={"source": _safe_rel(source, context.workspace), "target": _safe_rel(target, context.workspace), "commit": office_commit, "raw_bytes_hidden": True},
                        )
                if applied <= 0 and not allow_no_match:
                    if target.exists() and target.resolve() != source.resolve():
                        try:
                            target.unlink()
                        except OSError:
                            pass
                    return ToolResult(
                        invocation.step_id,
                        invocation.tool_name,
                        ToolResultStatus.FAILED,
                        "Office 文档未命中可替换文本，已停止写回。请先用 document_query 定位引用片段。",
                        error_code="document_rewrite_no_match",
                        data={"source": _safe_rel(source, context.workspace), "replacements": details, "raw_bytes_hidden": True},
                    )
        if dry_run:
            summary = "\n".join([
                "【文档写回预演】",
                f"- 源文件：{_safe_rel(source, context.workspace)}",
                f"- 目标文件：{_safe_rel(target, context.workspace)}",
                f"- 预计替换命中：{applied}",
                "- 状态：dry_run，仅经过治理链预演，未写入文件。",
            ])
            return ToolResult(invocation.step_id, invocation.tool_name, ToolResultStatus.OK, truncate_text(summary, context.policy.max_output_chars), data={"dry_run": True, "applied_count": applied, "replacements": details, "raw_bytes_hidden": True})

        if not target.exists() and applied > 0:
            return ToolResult(invocation.step_id, invocation.tool_name, ToolResultStatus.FAILED, "写回后目标文件未生成。", error_code="writeback_target_missing")
        new_hash = _sha256_file(target) if target.exists() else ""
        rollback_action = "restore_backup" if backup_rel else "delete_output"
        manifest = {
            "schema": WRITEBACK_SCHEMA,
            "operation_id": op_id,
            "created_at": datetime.now().isoformat(timespec="seconds"),
            "tool": "document_apply_rewrite",
            "context_source": context_source,
            "instruction": safe_text(instruction, 600),
            "source_path": _safe_rel(source, context.workspace),
            "target_path": _safe_rel(target, context.workspace),
            "backup_path": backup_rel,
            "rollback_action": rollback_action,
            "overwrite": bool(target.resolve() == source.resolve()),
            "operation": operation or "replace",
            "source_hash_before": source_hash,
            "target_hash_after": new_hash,
            "physical_commit_verified": True,
            "source_path_normalization": normalization_public_data(source_path_normalization),
            "target_path_normalization": normalization_public_data(target_path_normalization),
            "applied_count": applied,
            "replacement_summary": details,
            "encoding": encoding,
            "encoding_uncertain": encoding_uncertain,
            "raw_bytes_hidden": True,
            "quality_gate_chain": "RuntimeToolRegistry -> ExecutionSpine -> PermitGateway(A5Only) -> Adapter -> Audit",
        }
        manifest_rel = _write_manifest(context.workspace, op_id, manifest)
        lines = [
            "【文档写回】",
            f"- 状态：已{'覆盖写入' if target.resolve() == source.resolve() else '生成修订副本'}。",
            f"- 源文件：{_safe_rel(source, context.workspace)}",
            f"- 输出文件：{_safe_rel(target, context.workspace)}",
            f"- 替换命中：{applied}",
            f"- 操作 ID：{op_id}",
            f"- 回滚凭据：{manifest_rel}",
            "- 状态：已完成物理落盘验真。",
            "- 边界：写回只通过 Runtime 注册工具执行，已生成审计与回滚清单；主会话不展示原始正文或二进制。",
        ]
        artifacts = [_safe_rel(target, context.workspace), manifest_rel]
        if backup_rel:
            artifacts.append(backup_rel)
        return ToolResult(
            invocation.step_id,
            invocation.tool_name,
            ToolResultStatus.OK,
            truncate_text("\n".join(lines), context.policy.max_output_chars),
            data={
                "operation_id": op_id,
                "source": _safe_rel(source, context.workspace),
                "target": _safe_rel(target, context.workspace),
                "backup": backup_rel,
                "manifest": manifest_rel,
                "applied_count": applied,
                "rollback_action": rollback_action,
                "raw_bytes_hidden": True,
                "physical_commit_verified": True,
                "source_path_normalization": normalization_public_data(source_path_normalization),
                "target_path_normalization": normalization_public_data(target_path_normalization),
            },
            artifacts=artifacts,
        )
    except WorkspaceViolation as exc:
        return ToolResult(invocation.step_id, invocation.tool_name, ToolResultStatus.BLOCKED, str(exc), error_code="workspace_violation")
    except OSError as exc:
        return ToolResult(invocation.step_id, invocation.tool_name, ToolResultStatus.FAILED, f"文档写回失败：{exc}", error_code="os_error")
    except Exception as exc:
        return ToolResult(invocation.step_id, invocation.tool_name, ToolResultStatus.FAILED, f"文档写回失败：{safe_text(exc, 240)}", error_code="document_writeback_error")


def document_rollback_adapter(invocation: ToolInvocation, context: TurnContext) -> ToolResult:
    guard = WorkspaceGuard(context.workspace)
    try:
        args = dict(invocation.arguments)
        operation_id = safe_text(args.get("operation_id") or args.get("id"), 100)
        manifest_hint = args.get("manifest") or args.get("manifest_path") or args.get("path")
        if operation_id:
            manifest_path = _manifest_path(context.workspace, operation_id)
        elif manifest_hint:
            manifest_path = guard.resolve_for_read(manifest_hint)
        else:
            return ToolResult(invocation.step_id, invocation.tool_name, ToolResultStatus.FAILED, "缺少 operation_id 或 manifest_path，无法回滚。", error_code="rollback_manifest_missing")
        if not manifest_path.exists():
            return ToolResult(invocation.step_id, invocation.tool_name, ToolResultStatus.FAILED, "回滚清单不存在。", error_code="rollback_manifest_not_found")
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
        if manifest.get("schema") != WRITEBACK_SCHEMA:
            return ToolResult(invocation.step_id, invocation.tool_name, ToolResultStatus.FAILED, "回滚清单 schema 不匹配。", error_code="rollback_schema_mismatch")
        target_rel = safe_text(manifest.get("target_path"), 600)
        target = guard.resolve_for_write(target_rel)
        force = _bool(args.get("force"), default=False)
        expected_hash = safe_text(manifest.get("target_hash_after"), 128)
        if target.exists() and expected_hash and _sha256_file(target) != expected_hash and not force:
            return ToolResult(
                invocation.step_id,
                invocation.tool_name,
                ToolResultStatus.BLOCKED,
                "目标文件已在写回后发生变化。为避免覆盖后续编辑，已停止回滚；确认后可传 force=true。",
                error_code="rollback_target_changed",
                data={"target": target_rel, "operation_id": manifest.get("operation_id"), "raw_bytes_hidden": True},
            )
        op_id = safe_text(manifest.get("operation_id"), 100)
        pre_backup = _backup_file(context.workspace, target, op_id + "_rollback", "target_before_rollback") if target.exists() else ""
        action = safe_text(manifest.get("rollback_action") or ("restore_backup" if manifest.get("backup_path") else "delete_output"), 80)
        backup_rel = safe_text(manifest.get("backup_path"), 600)
        if action == "restore_backup":
            if not backup_rel:
                return ToolResult(invocation.step_id, invocation.tool_name, ToolResultStatus.FAILED, "回滚清单没有备份文件路径。", error_code="rollback_backup_missing")
            backup = guard.resolve_for_read(backup_rel)
            if not backup.exists():
                return ToolResult(invocation.step_id, invocation.tool_name, ToolResultStatus.FAILED, "备份文件不存在，无法回滚。", error_code="rollback_backup_not_found")
            target.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(backup, target)
            result_action = "已从备份恢复目标文件。"
        else:
            if target.exists():
                target.unlink()
            result_action = "已删除写回生成的新文件。"
        rollback_payload = {
            "schema": ROLLBACK_SCHEMA,
            "operation_id": op_id,
            "rolled_back_at": datetime.now().isoformat(timespec="seconds"),
            "target_path": target_rel,
            "backup_path": backup_rel,
            "pre_rollback_backup": pre_backup,
            "action": action,
            "force": force,
            "raw_bytes_hidden": True,
        }
        rollback_rel = _write_manifest(context.workspace, f"{op_id}_rollback", rollback_payload)
        summary = "\n".join([
            "【文档回滚】",
            f"- 状态：{result_action}",
            f"- 操作 ID：{op_id}",
            f"- 目标文件：{target_rel}",
            f"- 回滚记录：{rollback_rel}",
            "- 边界：回滚只基于写回清单和备份执行，未绕过 Runtime / QualityGate / Audit。",
        ])
        artifacts = [rollback_rel]
        if target.exists():
            artifacts.append(_safe_rel(target, context.workspace))
        if pre_backup:
            artifacts.append(pre_backup)
        return ToolResult(
            invocation.step_id,
            invocation.tool_name,
            ToolResultStatus.OK,
            truncate_text(summary, context.policy.max_output_chars),
            data={"operation_id": op_id, "target": target_rel, "rollback_manifest": rollback_rel, "pre_rollback_backup": pre_backup, "raw_bytes_hidden": True},
            artifacts=artifacts,
        )
    except WorkspaceViolation as exc:
        return ToolResult(invocation.step_id, invocation.tool_name, ToolResultStatus.BLOCKED, str(exc), error_code="workspace_violation")
    except OSError as exc:
        return ToolResult(invocation.step_id, invocation.tool_name, ToolResultStatus.FAILED, f"文档回滚失败：{exc}", error_code="os_error")
    except Exception as exc:
        return ToolResult(invocation.step_id, invocation.tool_name, ToolResultStatus.FAILED, f"文档回滚失败：{safe_text(exc, 240)}", error_code="document_rollback_error")
